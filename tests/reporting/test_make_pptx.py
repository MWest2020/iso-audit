"""Tests voor `iso_audit.reporting.make_pptx` — snapshot-presentatie."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from iso_audit.reporting import make_pptx


def test_maak_presentatie_genereert_bestand(tmp_path: Path) -> None:
    """Smoke test: alle 9 slides worden zonder crash gerenderd."""
    out = tmp_path / "snapshot.pptx"
    pad = make_pptx.maak_presentatie(uitvoer=out)
    assert pad == out
    assert out.is_file()
    assert out.stat().st_size > 1000


def test_maak_presentatie_aantal_slides(tmp_path: Path) -> None:
    """De snapshot heeft exact 9 slides (titel + 7 inhoud + conclusie)."""
    out = tmp_path / "snapshot.pptx"
    make_pptx.maak_presentatie(uitvoer=out)
    prs = Presentation(str(out))
    assert len(prs.slides) == 9


def test_maak_presentatie_maakt_parent_dir(tmp_path: Path) -> None:
    """Niet-bestaande parent-dir wordt aangemaakt."""
    out = tmp_path / "nested" / "dir" / "snapshot.pptx"
    pad = make_pptx.maak_presentatie(uitvoer=out)
    assert pad.is_file()


def test_slide_dimensies_widescreen(tmp_path: Path) -> None:
    """Slide-formaat is 13.33 x 7.5 inch (widescreen 16:9)."""
    out = tmp_path / "snapshot.pptx"
    make_pptx.maak_presentatie(uitvoer=out)
    prs = Presentation(str(out))
    assert prs.slide_width == make_pptx.SLIDE_W
    assert prs.slide_height == make_pptx.SLIDE_H


def test_individuele_slide_functies_werken() -> None:
    """Elke slide-functie kan zelfstandig draaien (geen onderling state-conflict)."""
    prs = Presentation()
    prs.slide_width = make_pptx.SLIDE_W
    prs.slide_height = make_pptx.SLIDE_H
    # Volgorde-onafhankelijk: 9 separate slide-builders.
    for slide_fn in (
        make_pptx.slide_titel,
        make_pptx.slide_organisatie,
        make_pptx.slide_aanpak,
        make_pptx.slide_cijfers,
        make_pptx.slide_kritiek,
        make_pptx.slide_plan_nc,
        make_pptx.slide_ofi,
        make_pptx.slide_positief,
        make_pptx.slide_conclusie,
    ):
        slide_fn(prs)
    assert len(prs.slides) == 9
